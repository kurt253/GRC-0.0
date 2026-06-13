"""Bouw een uitgebreide PPTX-presentatie van de GRC Tool."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── constanten ────────────────────────────────────────────────────────────────
OUT   = Path("data/GRC_Tool_Presentatie.pptx")
IMGS  = Path("data/presentation_screenshots")

NAVY  = RGBColor(0x0D, 0x2C, 0x6B)
BLUE  = RGBColor(0x1D, 0x4E, 0xD8)
LGRAY = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1E, 0x29, 0x3B)
GOLD  = RGBColor(0xF5, 0xA6, 0x23)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ── helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    return shp

def add_txt(slide, text, x, y, w, h, size=18, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def add_img(slide, path, x, y, w=None, h=None):
    path = str(path)
    if w and h:
        return slide.shapes.add_picture(path, x, y, w, h)
    elif w:
        return slide.shapes.add_picture(path, x, y, width=w)
    elif h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_txt(slide, title, Inches(0.35), Inches(0.12), Inches(11), Inches(0.6),
            size=26, bold=True, color=WHITE)
    if subtitle:
        add_txt(slide, subtitle, Inches(0.35), Inches(0.72), Inches(11), Inches(0.35),
                size=13, color=RGBColor(0xBA, 0xD3, 0xF8))

def footer(slide, text="GRC Tool  |  Informatieveiligheid voor Overheidsdiensten"):
    add_rect(slide, 0, SLIDE_H - Inches(0.28), SLIDE_W, Inches(0.28), NAVY)
    add_txt(slide, text, Inches(0.2), SLIDE_H - Inches(0.26),
            Inches(10), Inches(0.24), size=8,
            color=RGBColor(0xBA, 0xD3, 0xF8))

def bullet_box(slide, items, x, y, w, h, title=None, title_color=BLUE):
    add_rect(slide, x, y, w, h, LGRAY)
    ty = y
    if title:
        add_txt(slide, title, x + Inches(0.15), ty + Inches(0.1),
                w - Inches(0.2), Inches(0.35), size=12, bold=True, color=title_color)
        ty += Inches(0.38)
    for item in items:
        add_txt(slide, f"  {item}", x + Inches(0.05), ty,
                w - Inches(0.1), Inches(0.3), size=11, color=BLACK)
        ty += Inches(0.28)

def new_slide(title=None, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    if title:
        header_bar(s, title, subtitle)
    footer(s)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITEL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(0.18), SLIDE_H, BLUE)
add_rect(s, Inches(0.5), Inches(1.5), Inches(9.2), Inches(4.2), WHITE)

add_txt(s, "GRC Tool", Inches(0.75), Inches(1.75), Inches(8.5), Inches(1.1),
        size=48, bold=True, color=NAVY)
add_txt(s, "Informatieveiligheid voor Overheidsdiensten",
        Inches(0.75), Inches(2.8), Inches(8.5), Inches(0.65),
        size=20, color=BLUE)
add_txt(s, "Governance  |  Risicobeheer  |  Compliance",
        Inches(0.75), Inches(3.45), Inches(8.5), Inches(0.45),
        size=14, color=RGBColor(0x64, 0x74, 0x8B))
add_txt(s, "CyFun 2025  |  NIS2-conformiteit  |  Belgische overheid",
        Inches(0.75), Inches(4.8), Inches(8.5), Inches(0.4),
        size=12, color=WHITE)
add_txt(s, "v0.3  |  2026", Inches(10.5), Inches(6.9),
        Inches(2.5), Inches(0.35), size=10,
        color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Inhoud", "Overzicht van de presentatie")

agenda = [
    ("01", "Architectuur",               "Doel, technologiestack, dataflow"),
    ("02", "Excel-werkboek",             "Sheets, VBA-macro's, CIA-model"),
    ("03", "Configuratie",               "Taal, versie, assurance niveau"),
    ("04", "Processen & Assets",         "Processen, IA, afhankelijke assets"),
    ("05", "Maatregelen",                "CyFun controls + risicobeheerplan"),
    ("06", "Controles & Acties",         "Periodieke controles, actiebeheer"),
    ("07", "Kwetsbaarheden & Dashboard", "RARM-matrix, KPI-overzicht"),
    ("08", "Werkproces",                 "End-to-end van Access naar rapport"),
]

cw = Inches(3.1)
for i, (num, titel, desc) in enumerate(agenda):
    col = i % 4
    row = i // 4
    x = Inches(0.3) + col * (cw + Inches(0.18))
    y = Inches(1.35) + row * Inches(2.55)
    add_rect(s, x, y, cw, Inches(2.35), LGRAY)
    add_rect(s, x, y, Inches(0.5), Inches(2.35), BLUE)
    add_txt(s, num, x + Inches(0.07), y + Inches(0.9),
            Inches(0.38), Inches(0.45), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(s, titel, x + Inches(0.6), y + Inches(0.2),
            cw - Inches(0.7), Inches(0.5), size=12, bold=True, color=NAVY)
    add_txt(s, desc, x + Inches(0.6), y + Inches(0.75),
            cw - Inches(0.7), Inches(0.9), size=10, color=RGBColor(0x64, 0x74, 0x8B))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTUUR
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Architectuur", "Hoe werkt de GRC Tool?")

pijlers = [
    ("Access-database", ["MNMTool brondata", "Processen en assets", "Kwetsbaarheden",
                          "1-klik import via VBA", "ACE/Jet ADODB-connectie"], NAVY),
    ("Excel-werkboek (.xlsm)", ["Centrale datakern", "VBA-macro's (GRC_Macros)", "CyFun 2025 ingebouwd",
                                  "CIA-model per DA", "RARM-kwetsbaarhedenmatrix"], BLUE),
    ("Streamlit-dashboard", ["Leest en schrijft xlsm", "8 functionele pagina's", "Risicobeheerplan",
                               "Assurance-level filter", "Periodieke controles & acties"],
     RGBColor(0x05, 0x96, 0x69)),
]

for i, (titel, items, kleur) in enumerate(pijlers):
    x = Inches(0.3) + i * Inches(4.35)
    add_rect(s, x, Inches(1.3), Inches(4.1), Inches(0.42), kleur)
    add_txt(s, titel, x + Inches(0.1), Inches(1.32), Inches(3.9), Inches(0.38),
            size=12, bold=True, color=WHITE)
    for j, item in enumerate(items):
        yy = Inches(1.84) + j * Inches(0.84)
        add_rect(s, x, yy, Inches(4.1), Inches(0.76), LGRAY)
        add_rect(s, x, yy, Inches(0.06), Inches(0.76), kleur)
        add_txt(s, item, x + Inches(0.15), yy + Inches(0.23),
                Inches(3.8), Inches(0.42), size=11, color=BLACK)

add_txt(s, "<  Import  >", Inches(4.48), Inches(3.9), Inches(0.75), Inches(0.55),
        size=9, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_txt(s, "<  Lezen/Schrijven  >", Inches(8.65), Inches(3.9), Inches(1.1), Inches(0.55),
        size=9, bold=True, color=RGBColor(0x05, 0x96, 0x69), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — EXCEL SHEETS OVERZICHT
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Excel-werkboek", "GRC_Tool.xlsm — overzicht van alle tabbladen")

sheets = [
    ("Config",               "Organisatienaam, taal, versie, assurance niveau",    False),
    ("Info",                 "Gebruikershandleiding en snelkoppelingen",             False),
    ("Processes",            "Bedrijfsprocessen met CIA-classificatie (I en A)",     False),
    ("Information Assets",   "Informatieassets met vertrouwelijkheidsniveau (C)",    False),
    ("Dependent Assets",     "Afhankelijke assets + CIA-vereisten & objectieven",    False),
    ("RARM",                 "Risk Assessment & Remediation Matrix per DA",          False),
    ("Responsible Persons",  "Verantwoordelijken per domein",                        False),
    ("Import & Export",      "VBA-knoppen voor Access-import en export",             False),
    ("CyFun Controls",       "Alle CyFun 2025 controls (Basic/Important/Essential)", False),
    ("Kwetsbaarheden",       "Kwetsbaarhedenmatrix per afhankelijke asset",          False),
    ("Risicobeheer",         "Risicobeheerplan: status per maatregel per DA",        True),
    ("Controles",            "Periodieke controlesschema's per DA",                  True),
    ("Acties",               "Actielijst met opvolging en vervaldatums",             True),
]

cw2 = Inches(6.2)
for i, (naam, desc, is_new) in enumerate(sheets):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * (cw2 + Inches(0.4))
    y = Inches(1.3) + row * Inches(0.42)
    kleur = RGBColor(0x05, 0x96, 0x69) if is_new else NAVY
    add_rect(s, x, y, Inches(2.0), Inches(0.35), kleur)
    add_txt(s, naam, x + Inches(0.07), y + Inches(0.05),
            Inches(1.9), Inches(0.27), size=10, bold=True, color=WHITE)
    add_txt(s, desc, x + Inches(2.1), y + Inches(0.06),
            cw2 - Inches(2.15), Inches(0.27), size=10, color=BLACK)
    if is_new:
        add_txt(s, "NIEUW", x + Inches(2.1), y + Inches(0.06),
                Inches(0.6), Inches(0.27), size=9, bold=True,
                color=RGBColor(0x05, 0x96, 0x69))
        add_txt(s, desc, x + Inches(2.75), y + Inches(0.06),
                cw2 - Inches(2.8), Inches(0.27), size=10, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EXCEL CONFIG SCREENSHOT
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Excel — Config-tabblad", "Instellingen, assurance niveau en snelkoppelingen")
add_img(s, IMGS / "xl_01_config.png",
        Inches(0.3), Inches(1.2), h=Inches(5.9))
bullet_box(s, [
    "Naam overheidsinstelling en dienst/entiteit",
    "Taalinstelling: NL / FR / EN",
    "Dropdowns wisselen automatisch bij taalwijziging",
    "Versie, laatste update, bijgewerkt door",
    "Assurance niveau: Basic / Important / Essential",
    "Tip-sectie en snelkoppelingen naar alle tabs",
], Inches(5.1), Inches(1.25), Inches(7.95), Inches(3.2),
   title="Velden in Config-tabblad")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EXCEL IMPORT & EXPORT
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Excel — Import & Export", "VBA-macro's voor gegevensimport vanuit Access-database")
add_img(s, IMGS / "xl_10_import_export.png",
        Inches(0.3), Inches(1.2), h=Inches(5.9))
bullet_box(s, [
    "Importeer Alles: volledige dataset in 1 klik",
    "Importeer Informatieassets (T-Information Assets)",
    "Importeer Afhankelijke Assets + CIA-objectieven",
    "Importeer Processen + asset-koppelingen",
    "Importeer Kwetsbaarheden (vult RARM-matrix)",
    "Import Links DA/Kwetsbaarheden (na kwetsbaarheden)",
    "Tabelstructuur-referentie voor Access-ontwerpers",
], Inches(4.7), Inches(1.25), Inches(8.3), Inches(3.5),
   title="Importacties (VBA-macro's)")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — STREAMLIT CONFIGURATIE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Streamlit — Configuratie",
              "Assurance niveau instellen vanuit de Streamlit-app")
add_img(s, IMGS / "02_configuratie.png",
        Inches(0.25), Inches(1.2), w=Inches(8.5))
bullet_box(s, [
    "Toont Config-velden rechtstreeks uit xlsm",
    "Assurance niveau is verplicht in te stellen",
    "Rode waarschuwing als niveau niet is ingesteld",
    "Opslaan schrijft naar Config!D19 in xlsm",
    "Versie- en tijdstempelinfo weergegeven",
    "Maatregelen-pagina geblokkeerd zonder assurance niveau",
], Inches(9.0), Inches(1.25), Inches(4.05), Inches(3.0),
   title="Functionaliteit")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — STREAMLIT DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Streamlit — Dashboard",
              "Realtime KPI-overzicht en voortgang gegenereerd uit de xlsm")
add_img(s, IMGS / "01_dashboard.png",
        Inches(0.25), Inches(1.2), w=Inches(8.8))
bullet_box(s, [
    "9 processen | 8 informatieassets | 18 afhankelijke assets",
    "37 kwetsbaarheden | 34 controls (Basic-niveau)",
    "Taartdiagram: status risicobeheerplan",
    "Staafdiagram: controls per DA (assurance-gefilterd)",
    "Voortgangsbalk: uitgevoerde maatregelen per DA",
    "Open acties-tabel met vervaldatum",
], Inches(9.15), Inches(1.25), Inches(3.95), Inches(3.2),
   title="KPI-overzicht")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PROCESSEN & ASSETS (3 screenshots naast elkaar)
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Streamlit — Processen & Assets",
              "Drie tabbladen: Processen | Informatieassets | Afhankelijke Assets")

img_w = Inches(4.1)
tabs = [
    ("03_processen_tab1.png", "Processen",
     "CIA-labels (Laag/Gemiddeld/...) ipv nummers"),
    ("04_processen_tab2.png", "Informatieassets",
     "Vertrouwelijkheidsniveau als label"),
    ("05_processen_tab3.png", "Afhankelijke Assets",
     "CIA vereist vs. objectief + commentaar bewerkbaar"),
]
for i, (fname, tab, desc) in enumerate(tabs):
    x = Inches(0.25) + i * (img_w + Inches(0.28))
    add_rect(s, x, Inches(1.25), img_w, Inches(0.32), BLUE)
    add_txt(s, tab, x + Inches(0.1), Inches(1.27),
            img_w - Inches(0.1), Inches(0.28), size=11, bold=True, color=WHITE)
    add_img(s, IMGS / fname, x, Inches(1.57), w=img_w)
    add_txt(s, desc, x, Inches(5.2), img_w, Inches(0.65),
            size=9.5, color=RGBColor(0x64, 0x74, 0x8B))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — MAATREGELEN
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Streamlit — Maatregelen & Risicobeheerplan",
              "CyFun controls gefilterd op assurance niveau | CIA-gat detectie | Status per DA")
add_img(s, IMGS / "06_maatregelen.png",
        Inches(0.25), Inches(1.2), w=Inches(8.8))
bullet_box(s, [
    "Filter: Basic / Basic+Important / alle niveaus",
    "Per DA: CIA vereist vs. objectief (editeerbaar 1-5)",
    "Automatische CIA-gat detectie + waarschuwing",
    "Status: Uitgevoerd / Gepland / Niet uitvoeren",
    "Datum, verantwoordelijke en opmerkingen",
    "Voortgangsbalk per DA",
    "Opslaan: Risicobeheer-sheet + DA CIA-objectieven",
], Inches(9.15), Inches(1.25), Inches(3.95), Inches(4.0),
   title="Functionaliteit")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — PERIODIEKE CONTROLES
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Streamlit — Periodieke Controles",
              "Controlesschema per afhankelijke asset | Auto-actie bij uitvoering")
add_img(s, IMGS / "07_periodieke_controles.png",
        Inches(0.25), Inches(1.2), w=Inches(8.8))
bullet_box(s, [
    "Tabs per afhankelijke asset",
    "Frequentie: Maandelijks / Kwartaal / Jaarlijks / Ad hoc",
    "Controlemethode en artefactenlocatie vastleggen",
    "Laatste en volgende controledatum auto-berekend",
    "'Controle uitvoeren' maakt automatisch actie aan",
    "Status: OK / Te controleren / Vervallen",
    "Opslaan schrijft naar Controles-sheet in xlsm",
], Inches(9.15), Inches(1.25), Inches(3.95), Inches(3.8),
   title="Functionaliteit")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ACTIES
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Streamlit — Actiebeheer",
              "Gecentraliseerde actielijst met opvolging en rapportage")
add_img(s, IMGS / "08_acties.png",
        Inches(0.25), Inches(1.2), w=Inches(8.8))
bullet_box(s, [
    "KPI-balk: Open | In uitvoering | Gesloten | Vervallen",
    "Filter op status, DA, type en vervaldatum",
    "Actietypes: Periodieke controle / Maatregel / Opvolging",
    "Status inline bewerkbaar in tabel",
    "Nieuwe actie aanmaken via formulier",
    "Auto-increment ID: A-001, A-002, ...",
    "Opslaan schrijft naar Acties-sheet in xlsm",
], Inches(9.15), Inches(1.25), Inches(3.95), Inches(3.8),
   title="Functionaliteit")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — KWETSBAARHEDEN
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Streamlit — Kwetsbaarheden",
              "Kwetsbaarhedenmatrix en koppeling aan CyFun controls")
add_img(s, IMGS / "09_kwetsbaarheden.png",
        Inches(0.25), Inches(1.2), w=Inches(8.8))
bullet_box(s, [
    "Kwetsbaarheden geimporteerd via VBA-macro",
    "Matrix: koppeling kwetsbaarheid - CyFun control per DA",
    "Gefilterd op assurance niveau",
    "Inzicht in risico's per afhankelijke asset",
    "Basis voor risicobeheerplan en acties",
], Inches(9.15), Inches(1.25), Inches(3.95), Inches(2.5),
   title="Functionaliteit")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CIA-MODEL
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("CIA-model & Assurance Niveaus",
              "Classificatiemodel conform CyFun 2025")

cia = [
    ("C", "Vertrouwelijkheid", "Max. van gekoppelde informatieassets (IA)", BLUE),
    ("I", "Integriteit",       "Max. van gekoppelde bedrijfsprocessen",      NAVY),
    ("A", "Beschikbaarheid",   "Max. van gekoppelde bedrijfsprocessen",
     RGBColor(0x05, 0x96, 0x69)),
]
for i, (ltr, naam, bron, kleur) in enumerate(cia):
    x = Inches(0.3) + i * Inches(4.35)
    add_rect(s, x, Inches(1.35), Inches(4.1), Inches(1.75), LGRAY)
    add_rect(s, x, Inches(1.35), Inches(0.52), Inches(1.75), kleur)
    add_txt(s, ltr, x + Inches(0.1), Inches(1.7),
            Inches(0.34), Inches(0.62), size=26, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
    add_txt(s, naam, x + Inches(0.62), Inches(1.45),
            Inches(3.35), Inches(0.4), size=13, bold=True, color=kleur)
    add_txt(s, f"Vereist: {bron}",
            x + Inches(0.62), Inches(1.87), Inches(3.35), Inches(0.35),
            size=10, color=RGBColor(0x64, 0x74, 0x8B))
    add_txt(s, "Objectief: instelbaar in xlsm en Streamlit",
            x + Inches(0.62), Inches(2.2), Inches(3.35), Inches(0.35),
            size=10, color=BLACK)

niveaus = [("1 Laag", RGBColor(0x22,0xC5,0x5E)),
           ("2 Gemiddeld", RGBColor(0xEA,0xB3,0x08)),
           ("3 Hoog", RGBColor(0xF9,0x73,0x16)),
           ("4 Zeer Hoog", RGBColor(0xEF,0x44,0x44)),
           ("5 Kritiek", RGBColor(0x7C,0x3A,0xED))]
add_txt(s, "Classificatieschaal:", Inches(0.3), Inches(3.35),
        Inches(4), Inches(0.32), size=11, bold=True, color=NAVY)
for i, (niv, kl) in enumerate(niveaus):
    add_rect(s, Inches(0.3) + i * Inches(2.55), Inches(3.72), Inches(2.35), Inches(0.4), kl)
    add_txt(s, niv, Inches(0.35) + i * Inches(2.55), Inches(3.77),
            Inches(2.25), Inches(0.3), size=10, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)

add_txt(s, "Assurance-filter logica:", Inches(0.3), Inches(4.32),
        Inches(4), Inches(0.32), size=11, bold=True, color=NAVY)
assur = [("Basic",     "Toont enkel Basic controls", BLUE),
         ("Important", "Toont Basic + Important controls", NAVY),
         ("Essential", "Toont Basic + Important + Essential controls",
          RGBColor(0x7C,0x3A,0xED))]
for i, (niv, uitleg, kl) in enumerate(assur):
    x = Inches(0.3) + i * Inches(4.35)
    add_rect(s, x, Inches(4.7), Inches(4.1), Inches(1.35), LGRAY)
    add_rect(s, x, Inches(4.7), Inches(4.1), Inches(0.36), kl)
    add_txt(s, niv, x + Inches(0.1), Inches(4.72),
            Inches(3.9), Inches(0.32), size=12, bold=True, color=WHITE)
    add_txt(s, uitleg, x + Inches(0.1), Inches(5.1),
            Inches(3.9), Inches(0.65), size=11, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — WERKPROCES END-TO-END
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide("Werkproces end-to-end",
              "Van Access-database tot risicobeheerplan")

stappen = [
    ("1", "Access DB\nVoorbereiden",   "Vul processen, assets en kwetsbaarheden in de MNMTool-database"),
    ("2", "Importeer\nvia Excel",      "Klik 'Importeer Alles' op Import & Export-tab in Excel"),
    ("3", "Config\nInstellen",         "Stel assurance niveau in via Streamlit of Config-tab"),
    ("4", "RARM\nBekijken",            "RARM-tab synchroniseert DA-kolommen automatisch"),
    ("5", "Maatregelen\nOpvolgen",     "Status per control in Streamlit: Uitgevoerd/Gepland/Niet"),
    ("6", "Controles\nPlannen",        "Schema instellen per DA: frequentie, methode, artefacten"),
    ("7", "Acties\nOpvolgen",          "Werk acties bij; alles opgeslagen in xlsm"),
]

aw = Inches(1.72)
for i, (num, kort, lang) in enumerate(stappen):
    x = Inches(0.25) + i * (aw + Inches(0.08))
    kleur = BLUE if i % 2 == 0 else NAVY
    add_rect(s, x, Inches(1.35), aw, Inches(2.3), kleur)
    add_txt(s, num, x + Inches(0.1), Inches(1.42),
            Inches(0.38), Inches(0.4), size=18, bold=True, color=GOLD)
    add_txt(s, kort, x + Inches(0.1), Inches(1.82),
            aw - Inches(0.15), Inches(0.75), size=11, bold=True, color=WHITE)
    if i < 6:
        add_txt(s, ">", x + aw + Inches(0.0), Inches(2.1),
                Inches(0.12), Inches(0.4), size=16, bold=True, color=GOLD)
    add_txt(s, lang, x, Inches(3.82),
            aw, Inches(1.5), size=9.5, color=RGBColor(0x64, 0x74, 0x8B))

# Slotlijn
add_rect(s, Inches(0.25), Inches(5.55), Inches(12.85), Inches(0.04),
         RGBColor(0xE2, 0xE8, 0xF0))
add_txt(s, "Alle wijzigingen worden rechtstreeks weggeschreven naar GRC_Tool.xlsm (keep_vba=True) — macro's blijven intact.",
        Inches(0.3), Inches(5.65), Inches(12.5), Inches(0.4),
        size=10, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — SLOTPAGINA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(0.18), SLIDE_H, GOLD)
add_rect(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(3.2),
         RGBColor(0x13, 0x3A, 0x80))

add_txt(s, "GRC Tool v0.3", Inches(0.5), Inches(2.3),
        Inches(12.3), Inches(1.1), size=42, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
add_txt(s, "Informatieveiligheid voor Overheidsdiensten",
        Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
        size=18, color=RGBColor(0xBA, 0xD3, 0xF8), align=PP_ALIGN.CENTER)
add_txt(s, "CyFun 2025  |  NIS2  |  Excel + Streamlit  |  Belgische overheid",
        Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.45),
        size=13, color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)


# ── opslaan ───────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Presentatie opgeslagen: {OUT}  ({len(prs.slides)} slides)")
